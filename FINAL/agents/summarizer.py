import os
import json
from docx import Document
from pptx import Presentation
from llm.llama_wrapper import query_llama
import pdfplumber

# 🧠 Extract text from different file types
def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[-1].lower()

    if ext == '.docx':
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    elif ext == '.txt':
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    elif ext == '.json':
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return json.dumps(data, indent=2)

    elif ext == '.pptx':
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    
    elif ext == '.pdf':
        with pdfplumber.open(file_path) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)

    else:
        raise ValueError(f"Unsupported file type: {ext}")

# 🧠 Summarize individual text
def summarize_text(text: str, max_tokens: int = 1024) -> str:
    system_instruction = (
        "You are a helpful assistant. Summarize the following document into a concise summary:"
    )
    prompt = f"{system_instruction}\n\n{text[:max_tokens * 4]}\n\nSummary:"
    return query_llama(prompt, stream=False)

# 🧠 Summarize all files separately with labels
def summarize_files(file_paths: list[str]) -> str:
    all_summaries = []

    for path in file_paths:
        file_name = os.path.basename(path)
        try:
            text = extract_text_from_file(path)
            summary = summarize_text(text)
            block = f"\n📄 **Summary for:** `{file_name}`\n\n{summary.strip()}\n"
            all_summaries.append(block)
        except Exception as e:
            all_summaries.append(f"\n❌ Error summarizing `{file_name}`: {e}\n")

    return "\n".join(all_summaries)
